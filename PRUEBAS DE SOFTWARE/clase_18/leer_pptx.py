from pptx import Presentation

prs = Presentation('Sesion_18_Binomiales_LCS.pptx')
for i, slide in enumerate(prs.slides):
    texts = [s.text.strip() for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]
    if texts:
        print(f'=== Diapositiva {i+1} ===')
        for t in texts:
            print(t)
        print()
